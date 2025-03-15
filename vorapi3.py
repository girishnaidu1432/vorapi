from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.responses import JSONResponse, FileResponse
import openai
import PyPDF2
import pandas as pd
import pptx
from io import BytesIO
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import uvicorn
from fastapi.responses import StreamingResponse

app = FastAPI()

# OpenAI API Configuration
openai.api_key = "14560021aaf84772835d76246b53397a"
openai.api_base = "https://amrxgenai.openai.azure.com/"
openai.api_type = "azure"
openai.api_version = "2024-02-15-preview"
deployment_name = "gpt"

session_state = {
    "messages": [],
    "uploaded_file": None,
    "extracted_text": "",
}

# Function to extract text from uploaded files
def extract_text(file: UploadFile):
    text = ""
    content = file.file.read()
    file.file.seek(0)
    
    if file.content_type == "application/pdf":
        reader = PyPDF2.PdfReader(BytesIO(content))
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    
    elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(BytesIO(content))
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = pptx.Presentation(BytesIO(content))
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file.content_type == "text/csv":
        df = pd.read_csv(BytesIO(content))
        text += df.to_string()

    elif file.content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        df = pd.read_excel(BytesIO(content))
        text += df.to_string()

    return text.strip()

@app.get("/")
async def home():
    return {"message": "Welcome to the FastAPI Chatbot"}

@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    session_state["uploaded_file"] = file.filename
    session_state["extracted_text"] = extract_text(file)

    session_state["messages"].insert(0, {
        "role": "system",
        "content": f"Document Context:\n\n{session_state['extracted_text']}"
    })
    
    return {"filename": file.filename, "extracted_text": session_state["extracted_text"]}

@app.post("/chat/")
async def chat(user_input: str = Form(...)):
    if not user_input.strip():
        raise HTTPException(status_code=400, detail="User input cannot be empty")

    combined_prompt = user_input
    if session_state["uploaded_file"]:
        combined_prompt = f"Here is a document that provides context:\n\n{session_state['extracted_text']}\n\nNow, based on this document, answer the following:\n{user_input}"

    session_state["messages"].append({"role": "user", "content": user_input, "file": session_state["uploaded_file"]})

    response = openai.ChatCompletion.create(
        engine=deployment_name,
        messages=session_state["messages"],
        temperature=0.7,
        max_tokens=2000
    )

    ai_response = response["choices"][0]["message"]["content"]
    session_state["messages"].append({"role": "assistant", "content": ai_response})

    return JSONResponse(content={"response": ai_response})

# Function to create a DOCX file from chat history
def create_docx():
    doc = Document()
    doc.add_heading("Chatbot Conversation History", level=1).alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    for msg in session_state["messages"]:
        p = doc.add_paragraph()
        run = p.add_run(msg["role"].capitalize() + "\n")
        run.bold = True
        run.font.size = Pt(14)
        if "file" in msg and msg["file"]:
            doc.add_paragraph(f"📄 {msg['file']}")
        p.add_run(msg["content"]).font.size = Pt(12)
        p.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
        doc.add_paragraph("\n")

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


from fastapi.responses import StreamingResponse



@app.get("/download/")
async def download_chat():
    if not session_state["messages"]:
        raise HTTPException(status_code=400, detail="No chat history to download")

    docx_file = create_docx()  # Generate DOCX in memory
    docx_file.seek(0)  # Ensure file pointer is at the start

    return StreamingResponse(
        docx_file,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={
            "Content-Disposition": "attachment; filename=chat_history.docx",
            "Content-Type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        }
    )


@app.delete("/clear_chat/")
async def clear_chat():
    session_state["messages"] = []  # Reset chat history
    return {"message": "Chat history cleared successfully!"}


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
